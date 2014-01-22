# -*- coding: utf-8 -*-

# Copyright 2013 by Fabian Hachenberg http://github.com/fHachenberg
#
# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU General Public License for more details.
#
# You should have received a copy of the GNU General Public License
# along with this program.  If not, see <http://www.gnu.org/licenses/>.

'''
After a round has been created in Tournaman, this tool can read the round definitions from the Tournaman data files in the "Data" subdirectory (xml and dat) and then create the round slides in pptx format by substituting the respective entries in the presentation template. Has be called in the folder, in which the trm file resides (contains subdirectories "Data", "Registration", ...).

Normally, the only necessary input is the round_no. All other filenames are deducted from this. If you're using another branch in your tournament (The predefined one is "main") you need to specifiy this as well. Nevertheless all input and output filenames can also be set manually.
'''

__author__ = "Fabian Hachenberg"
__copyright__ = "Copyright 2013, Fabian Hachenberg"
__credits__ = ["Sönke Senff", "Jens Henning Fischer"]
__license__ = "GPL"
__version__ = "3"
__maintainer__ = "Fabian Hachenberg"
__email__ = "https://github.com/fHachenberg"
__status__ = "Production"

import argparse
import os
import re

from pptx import Presentation

import tournaman

parser = argparse.ArgumentParser(description=__doc__)
parser.add_argument('--round_no', dest='round_no', action='store', type=int, help='Index of round (1-N)', required=True)
parser.add_argument('--branch', dest='branch', action="store", type=str, default="main", help="Tournament branch")
parser.add_argument('--template', dest='pptx_template', action='store', help='Presentation template in pptx format', default="round presentation template.pptx")
parser.add_argument('--output', dest='output', action='store', default=None, help='Output filename')

parser.add_argument('--venue_file', dest='venue_file', action='store',
                   help='Manual choice of venue definition file (.def)', default=None)
parser.add_argument('--team_file', dest='team_file', action='store',
                   help='Manual choice of team definition file (.xml)', default=None)
parser.add_argument('--adjudicator_file', dest='adjudicator_file', action='store',
                   help='Manual choice of adjudicator definition file (.xml)', default=None)
parser.add_argument('--round_file', dest='round_file', action='store',
                   help='Manual choice of debate definition file (.xml)', default=None)

args = parser.parse_args()

round_no   = args.round_no
branch     = args.branch

#if manual data files are not specified, they are deducted automatically from round_no and branch
if args.venue_file == None:
    venue_file = os.path.join("Data", "venues%d-%s.dat" % (round_no, branch))
else:
    venue_file = args.venue_file
if args.team_file == None:
    team_file  = os.path.join("Data", "teams%d-%s.xml" % (round_no, branch))
else:
    team_file  = args.team_file
if args.adjudicator_file == None:
    adjudicator_file = os.path.join("Data", "adjudicators%d-%s.xml" % (round_no, branch))
else:
    adjudicator_file = args.adjudicator_file
if args.round_file == None:
    round_file = os.path.join("Data", "debates%d-%s.xml" % (round_no, branch))
else:
    round_file = args.round_file

#if output filename is not specified, it is generated from round_no and branch
if args.output == None:
    output = 'round presentation %s round %d.pptx' % (branch, round_no)
else:
    output = args.output

pptx_template = args.pptx_template

#Read data from tournaman files
venue_db = tournaman.parse_venue_def(venue_file)
team_db  = tournaman.parse_team_xml(team_file)
adjud_db = tournaman.parse_adjudicator_xml(adjudicator_file)
rnd = tournaman.parse_debates_xml(round_file, team_db, adjud_db, venue_db)

prs = Presentation(pptx_template)

replacements = [
                ("m", rnd.motion),
                ("r", str(round_no))
               ]
for i, debate in enumerate(rnd.debates):
    replacements += [("v%d"  % (i+1), debate.venue.name),
                     ("og%d" % (i+1), debate.og.name),
                     ("oo%d" % (i+1), debate.oo.name),
                     ("cg%d" % (i+1), debate.cg.name),
                     ("co%d" % (i+1), debate.co.name),
                     ("j%d"  % (i+1), "\n".join(adjud.name for adjud in debate.adjuds))]
replacements = dict(replacements)

valid_placeholder_patterns = r"m|v[0-9]+|j[0-9]+|og[0-9]+|cg[0-9]+|oo[0-9]+|co[0-9]+"

def do_repl(match):
    '''
    this callable does the replacements
    If a placeholder is not in the replacements dict,
    it is checked whether the placeholder conforms to one
    of the valid placeholder patterns. If not, an error is prompted.
    Otherwise a warning
    '''
    name = match.groups(1)[0]
    if name not in replacements.keys():
        if not re.match(valid_placeholder_patterns, name):
            print "ERROR -- invalid placeholder found in presentation:", name
        else:
            print "WARNING -- no data available to substitute valid placeholder'" + name + "'"
        return ""
    return replacements[name]

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_textframe:
            continue
        for paragraph in shape.textframe.paragraphs:
            for run in paragraph.runs:     
                if run.text == None:
                    continue
                run.text = re.sub(r"#([a-z0-9]+)", do_repl, run.text)

prs.save(output)

